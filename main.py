"""
main.py — PowerPoint 슬라이드 번역기 CLI 진입점

사용법:
    python main.py <input.pptx> <target_lang> [--output output.pptx]

예시:
    python main.py presentation.pptx ko
    python main.py presentation.pptx ja --output translated_ja.pptx
"""

import argparse
import copy
import logging
import os
import shutil
import sys

from dotenv import load_dotenv
from pptx import Presentation
from tqdm import tqdm

from pptx_handler import (
    extract_slide_context,
    extract_styled_paragraphs,
    apply_translated_runs,
    backup_paragraph_xml,
    iter_translatable_shapes,
    get_target_font,
)
from translator import (
    get_presentation_summary,
    translate_styled_text,
    translate_slide_batch,
    get_lang_name,
)

# ──────────────────────────────────────────────
#  로깅 설정
# ──────────────────────────────────────────────

logging.basicConfig(
    level=logging.WARNING,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────
#  출력 파일명 생성
# ──────────────────────────────────────────────

def make_output_path(input_path: str, target_lang: str) -> str:
    """입력 파일 경로에서 _언어코드.pptx 형태의 출력 경로를 생성합니다."""
    base, ext = os.path.splitext(input_path)
    return f"{base}_{target_lang}{ext}"


def _parse_slide_range(slides_arg: str, total_slides: int) -> set[int]:
    """
    슬라이드 범위 문자열을 파싱합니다.

    지원 형식:
        "5"      → {5}
        "3-10"   → {3, 4, 5, 6, 7, 8, 9, 10}
        "1-5"    → {1, 2, 3, 4, 5}

    Args:
        slides_arg: 사용자 입력 문자열
        total_slides: 전체 슬라이드 수 (범위 검증용)

    Returns:
        번역 대상 슬라이드 번호 집합 (1-based)
    """
    slides_arg = slides_arg.strip()

    if "-" in slides_arg:
        parts = slides_arg.split("-", 1)
        try:
            start = int(parts[0].strip())
            end = int(parts[1].strip())
        except ValueError:
            logger.error(f"슬라이드 범위 형식 오류: '{slides_arg}' (예: 3-10)")
            sys.exit(1)
    else:
        try:
            start = end = int(slides_arg)
        except ValueError:
            logger.error(f"슬라이드 번호 형식 오류: '{slides_arg}' (예: 5 또는 3-10)")
            sys.exit(1)

    # 범위 검증
    start = max(1, start)
    end = min(total_slides, end)
    if start > end:
        logger.error(f"슬라이드 범위 오류: {start}~{end} (전체 {total_slides}장)")
        sys.exit(1)

    return set(range(start, end + 1))


# ──────────────────────────────────────────────
#  텍스트 프레임 번역
# ──────────────────────────────────────────────

def translate_text_frame(text_frame, context: str, target_lang: str,
                         target_font: str | None = None,
                         pres_summary: str = "") -> bool:
    """
    텍스트 프레임 내 텍스트를 번역합니다.

    Returns:
        True: 번역 성공, False: 스킵 또는 실패
    """
    styled_data = extract_styled_paragraphs(text_frame)
    if styled_data is None:
        return False  # 텍스트 없음

    # 원본 텍스트 추출 (로그용)
    original_text = " | ".join(
        "".join(r["text"] for r in p["runs"])
        for p in styled_data["paragraphs"]
    ).strip()

    # GPT 번역 요청
    translated = translate_styled_text(styled_data, context, target_lang,
                                       pres_summary=pres_summary)
    if translated is None:
        logger.warning("  번역 실패 — 원문 유지")
        return False

    # 번역 결과 적용 (run.text 교체 방식 — XML 조작 없음)
    try:
        apply_translated_runs(text_frame, translated, styled_data["styles"],
                              target_font=target_font, target_lang=target_lang,
                              rPr_xml_map=styled_data.get("rPr_xml_map"))
        # 교체 로그
        translated_text = " | ".join(
            "".join(r["text"] for r in p["runs"])
            for p in translated["paragraphs"]
        ).strip()
        logger.info(f"  [교체 완료] \"{original_text}\" → \"{translated_text}\"")
        return True
    except Exception as e:
        logger.error(f"  번역 적용 중 오류: {e}")
        return False


# ──────────────────────────────────────────────
#  테이블 번역
# ──────────────────────────────────────────────

def translate_table(table, context: str, target_lang: str,
                    target_font: str | None = None,
                    pres_summary: str = "") -> int:
    """
    테이블의 각 셀을 번역합니다.

    Returns:
        번역된 셀 수
    """
    translated_count = 0
    for row in table.rows:
        for cell in row.cells:
            if cell.text.strip():
                success = translate_text_frame(cell.text_frame, context, target_lang,
                                              target_font, pres_summary=pres_summary)
                if success:
                    translated_count += 1
    return translated_count


# ──────────────────────────────────────────────
#  슬라이드 번역
# ──────────────────────────────────────────────

def translate_slide(slide, slide_num: int, target_lang: str,
                    pres_summary: str = "",
                    recent_translations: list[dict] | None = None) -> tuple[dict, list[dict]]:
    """
    슬라이드 하나를 번역합니다.
    텍스트 프레임 + 테이블 셀을 모두 수집 → 1회 API 호출로 일괄 번역 → 결과 적용.
    Phase 0 프레젠테이션 요약 + 직전 슬라이드 번역 이력을 맥락으로 활용합니다.

    Returns:
        (stats, slide_pairs)
        stats: {"text_frames": N, "tables": N, "cells": N} 번역 통계
        slide_pairs: 이 슬라이드의 원문→번역 쌍 [{"src": ..., "tgt": ...}, ...]
    """
    stats = {"text_frames": 0, "tables": 0, "cells": 0, "notes": 0}
    slide_pairs: list[dict] = []
    target_font = get_target_font(target_lang)

    # 슬라이드에 텍스트가 있는지 빠르게 확인
    slide_text = extract_slide_context(slide)
    if not slide_text.strip():
        logger.info(f"  슬라이드 {slide_num}: 텍스트 없음 — 스킵")
        return stats, slide_pairs

    # 텍스트 프레임 + 테이블 셀 + 노트를 하나의 배치로 수집
    batch_items = []   # (text_frame, styled_data, box_id, source)
    box_counter = 0
    table_cell_ids = {}  # box_id → table index (통계용)

    for shape, shape_type in iter_translatable_shapes(slide):
        if shape_type == "text_frame":
            styled_data = extract_styled_paragraphs(shape.text_frame)
            if styled_data is not None:
                box_id = f"T{box_counter}"
                box_counter += 1
                batch_items.append((shape.text_frame, styled_data, box_id, "text_frame"))
        elif shape_type == "table":
            table_idx = stats["tables"] + len(table_cell_ids) + 1
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        styled_data = extract_styled_paragraphs(cell.text_frame)
                        if styled_data is not None:
                            box_id = f"C{box_counter}"
                            box_counter += 1
                            batch_items.append((cell.text_frame, styled_data, box_id, "table_cell"))
                            table_cell_ids[box_id] = table_idx

    # 슬라이드 노트
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf.text.strip():
                styled_data = extract_styled_paragraphs(notes_tf)
                if styled_data is not None:
                    box_id = f"N{box_counter}"
                    box_counter += 1
                    batch_items.append((notes_tf, styled_data, box_id, "note"))
    except Exception:
        pass  # 노트 없거나 접근 불가 시 무시

    if not batch_items:
        logger.info(f"  슬라이드 {slide_num}: 번역할 항목 없음 — 스킵")
        return stats

    # 일괄 번역 (1회 API 호출)
    batch_input = [
        {"box_id": box_id, "styled_data": sd}
        for (_, sd, box_id, _) in batch_items
    ]
    logger.info(f"  슬라이드 {slide_num}: {len(batch_input)}개 항목 일괄 번역 요청...")
    translated_map = translate_slide_batch(
        batch_input, "", target_lang, pres_summary=pres_summary,
        recent_translations=recent_translations,
    )

    if translated_map is None:
        logger.warning(f"  슬라이드 {slide_num}: 일괄 번역 실패 — 개별 번역으로 폴백")
        for text_frame, styled_data, box_id, source in batch_items:
            translated = translate_styled_text(
                styled_data, "", target_lang, pres_summary=pres_summary
            )
            if translated is not None:
                try:
                    apply_translated_runs(
                        text_frame, translated, styled_data["styles"],
                        target_font=target_font, target_lang=target_lang,
                        rPr_xml_map=styled_data.get("rPr_xml_map")
                    )
                    if source == "text_frame":
                        stats["text_frames"] += 1
                    elif source == "note":
                        stats["notes"] += 1
                    else:
                        stats["cells"] += 1
                except Exception as e:
                    logger.error(f"  [{box_id}] 번역 적용 중 오류: {e}")
    else:
        for text_frame, styled_data, box_id, source in batch_items:
            if box_id not in translated_map:
                logger.warning(f"  [{box_id}] 번역 결과 누락 — 원문 유지")
                continue
            try:
                apply_translated_runs(
                    text_frame, translated_map[box_id], styled_data["styles"],
                    target_font=target_font, target_lang=target_lang,
                    rPr_xml_map=styled_data.get("rPr_xml_map")
                )
                if source == "text_frame":
                    stats["text_frames"] += 1
                elif source == "note":
                    stats["notes"] += 1
                else:
                    stats["cells"] += 1
            except Exception as e:
                logger.error(f"  [{box_id}] 번역 적용 중 오류: {e}")

    # 테이블 수 집계
    if table_cell_ids:
        stats["tables"] = len(set(table_cell_ids.values()))

    # 원문→번역 쌍 수집 (용어 일관성 참조용)
    if translated_map:
        for _, styled_data, box_id, _ in batch_items:
            if box_id not in translated_map:
                continue
            src_text = " ".join(
                r["text"] for p in styled_data["paragraphs"] for r in p["runs"]
            ).strip()
            tgt_text = " ".join(
                r.get("text", "") for p in translated_map[box_id].get("paragraphs", [])
                for r in p.get("runs", [])
            ).strip()
            if src_text and tgt_text and src_text != tgt_text:
                slide_pairs.append({"src": src_text, "tgt": tgt_text})

    return stats, slide_pairs


# ──────────────────────────────────────────────
#  메인 파이프라인
# ──────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="PowerPoint 슬라이드를 지정 언어로 번역합니다.",
        epilog="예시: python main.py presentation.pptx ko\n"
               "      python main.py presentation.pptx ko --slides 5-10\n"
               "      python main.py presentation.pptx ko --slides 3",
    )
    parser.add_argument(
        "input_file",
        help="번역할 PPTX 파일 경로",
    )
    parser.add_argument(
        "target_lang",
        help="대상 언어 코드 (예: ko, ja, zh, en, es, fr, de)",
    )
    parser.add_argument(
        "--output", "-o",
        help="출력 파일 경로 (미지정 시 원본명_언어코드.pptx)",
        default=None,
    )
    parser.add_argument(
        "--slides", "-s",
        help="번역할 슬라이드 범위 (예: 5, 3-10, 1-5). 미지정 시 전체 번역",
        default=None,
    )
    parser.add_argument(
        "--verbose", "-v",
        help="상세 로그 출력 (기본: 생략)",
        action="store_true",
        default=False,
    )
    args = parser.parse_args()

    # ── 로그 레벨 설정 ──
    if args.verbose:
        logging.getLogger().setLevel(logging.INFO)
        logging.getLogger("translator").setLevel(logging.INFO)
        logging.getLogger("pptx_handler").setLevel(logging.INFO)

    # ── 환경변수 로드 ──
    load_dotenv()

    # ── 입력 파일 검증 ──
    if not os.path.isfile(args.input_file):
        logger.error(f"파일을 찾을 수 없습니다: {args.input_file}")
        sys.exit(1)

    if not args.input_file.lower().endswith(".pptx"):
        logger.error("PPTX 파일만 지원합니다.")
        sys.exit(1)

    # ── 출력 파일 경로 결정 ──
    output_path = args.output or make_output_path(args.input_file, args.target_lang)

    # ── 파일 복제 ──
    logger.info(f"파일 복제: {args.input_file} → {output_path}")
    shutil.copy2(args.input_file, output_path)

    # ── PPTX 열기 ──
    logger.info(f"PPTX 파일 로드 중: {output_path}")
    prs = Presentation(output_path)

    total_slides = len(prs.slides)
    lang_name = get_lang_name(args.target_lang)

    # ── 슬라이드 범위 파싱 ──
    slide_range = None  # None = 전체
    if args.slides:
        slide_range = _parse_slide_range(args.slides, total_slides)
        if len(slide_range) > 1:
            range_desc = f"{min(slide_range)}-{max(slide_range)}장"
        else:
            range_desc = f"{min(slide_range)}장"
        logger.info(f"번역 대상: {total_slides}개 슬라이드 중 {range_desc} ({len(slide_range)}개) → {lang_name}")
    else:
        logger.info(f"번역 시작: {total_slides}개 슬라이드 전체 → {lang_name}")
    logger.info("=" * 60)

    target_count = len(slide_range) if slide_range else total_slides

    # ── Phase 0: 프레젠테이션 전체 맥락 파악 (상위 5장, 1회 API) ──
    pres_summary = ""
    top_n = min(5, total_slides)
    logger.info(f"프레젠테이션 전체 맥락 파악 중 (상위 {top_n}장 분석)...")
    top_texts = []
    for slide in list(prs.slides)[:top_n]:
        slide_text = extract_slide_context(slide)
        if slide_text.strip():
            top_texts.append(slide_text)
    if top_texts:
        combined = "\n---\n".join(top_texts)
        pres_summary = get_presentation_summary(combined, args.target_lang)
        logger.info(f"프레젠테이션 요약: {pres_summary[:120]}...")
    else:
        logger.info("텍스트 없음 — 맥락 요약 생략")

    # ── 슬라이드별 번역 ──
    total_stats = {"text_frames": 0, "tables": 0, "cells": 0, "notes": 0}
    # 직전 3장 번역 이력 (용어 일관성 유지용 슬라이딩 윈도우)
    recent_history: list[list[dict]] = []  # [slide_pairs, slide_pairs, ...]

    with tqdm(total=target_count, desc="번역 진행", unit="slide") as pbar:
        for slide_idx, slide in enumerate(prs.slides, 1):
            # 범위 지정 시 해당 슬라이드만 번역
            if slide_range and slide_idx not in slide_range:
                continue

            # 직전 3장의 번역 쌍을 평탄화하여 전달
            flat_recent = [pair for pairs in recent_history for pair in pairs]

            stats, slide_pairs = translate_slide(
                slide, slide_idx, args.target_lang,
                pres_summary=pres_summary,
                recent_translations=flat_recent if flat_recent else None,
            )

            # 슬라이딩 윈도우 갱신 (최근 3장 유지)
            if slide_pairs:
                recent_history.append(slide_pairs)
                if len(recent_history) > 3:
                    recent_history.pop(0)

            total_stats["text_frames"] += stats["text_frames"]
            total_stats["tables"] += stats["tables"]
            total_stats["cells"] += stats["cells"]
            total_stats["notes"] += stats["notes"]

            pbar.update(1)
            pbar.set_postfix(
                텍스트=total_stats["text_frames"],
                표=total_stats["tables"],
                노트=total_stats["notes"],
            )

    # ── 저장 ──
    logger.info("=" * 60)
    logger.info("번역 완료. 파일 저장 중...")
    prs.save(output_path)

    # ── 결과 출력 ──
    logger.info(f"✓ 출력 파일: {output_path}")
    logger.info(
        f"✓ 통계: 텍스트 프레임 {total_stats['text_frames']}개, "
        f"표 {total_stats['tables']}개 ({total_stats['cells']}셀), "
        f"노트 {total_stats['notes']}개 번역"
    )


if __name__ == "__main__":
    main()
