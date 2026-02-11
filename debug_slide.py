"""문제 슬라이드 분석: Microsoft is built on trust"""
from pptx import Presentation
from pptx_handler import extract_styled_paragraphs
from lxml import etree

NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

# 원본과 번역본 모두 분석
for fname in ['Foundry_L300.PPTX', 'Foundry_L300_ko.PPTX']:
    prs = Presentation(fname)
    for si, slide in enumerate(prs.slides):
        full_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text += shape.text_frame.text + " "
        
        if 'trust' in full_text.lower() or '신뢰' in full_text:
            if 'data' in full_text.lower() or '데이터' in full_text:
                print(f"\n{'='*60}")
                print(f"파일: {fname}, 슬라이드 {si+1}")
                print(f"{'='*60}")
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    tf = shape.text_frame
                    text = tf.text.strip()
                    if not text:
                        continue
                    print(f"\n  Shape {shape.shape_id}: '{text[:60]}...' " if len(text)>60 else f"\n  Shape {shape.shape_id}: '{text}'")
                    
                    for pi, para in enumerate(tf.paragraphs):
                        runs = list(para.runs)
                        if not runs:
                            continue
                        for ri, run in enumerate(runs):
                            rPr = run._r.find(f'{NS}rPr')
                            attrs = {}
                            if rPr is not None:
                                attrs['sz'] = rPr.get('sz')
                                attrs['b'] = rPr.get('b')
                                attrs['i'] = rPr.get('i')
                                # color
                                solid = rPr.find(f'{NS}solidFill')
                                if solid is not None:
                                    srgb = solid.find(f'{NS}srgbClr')
                                    if srgb is not None:
                                        attrs['color'] = srgb.get('val')
                                grad = rPr.find(f'{NS}gradFill')
                                if grad is not None:
                                    attrs['gradFill'] = True
                            print(f"    P{pi}R{ri}: '{run.text}' | {attrs}")
                break  # 한 슬라이드만
