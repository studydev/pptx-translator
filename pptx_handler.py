"""
pptx_handler.py — PPTX 파싱, 스타일 추출/적용 모듈

슬라이드 내 텍스트 프레임과 테이블에서 Run 수준 스타일을 추출하고,
번역 결과를 스타일을 보존하며 다시 적용합니다.

핵심 원칙: XML 직접 조작 금지 — run.text 교체만으로 서식 100% 보존
"""

import copy
import logging
from typing import Any
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────
#  스타일 추출 / 비교 헬퍼
# ──────────────────────────────────────────────

def _extract_run_style(run) -> dict:
    """Run의 서식 속성을 딕셔너리로 추출합니다."""
    font = run.font
    style: dict[str, Any] = {}

    # 기본 속성
    if font.name is not None:
        style["name"] = font.name
    if font.size is not None:
        style["size"] = font.size  # EMU 단위 정수
    if font.bold is not None:
        style["bold"] = font.bold
    if font.italic is not None:
        style["italic"] = font.italic
    if font.underline is not None:
        style["underline"] = font.underline

    # 색상 — XML에서 직접 읽기 (python-pptx font.color 접근 시 gradFill이 파괴됨)
    try:
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            solid = rPr.find(qn('a:solidFill'))
            grad = rPr.find(qn('a:gradFill'))
            if solid is not None:
                srgb = solid.find(qn('a:srgbClr'))
                if srgb is not None:
                    style["color_rgb"] = srgb.get('val')
                    style["color_type"] = "solid"
            elif grad is not None:
                # 그라데이션은 첫번째 색상만 참고용으로 추출
                gs_list = grad.findall(f'.//{qn("a:gs")}')
                if gs_list:
                    srgb = gs_list[0].find(qn('a:srgbClr'))
                    if srgb is not None:
                        style["color_rgb"] = srgb.get('val')
                style["color_type"] = "gradient"
    except (AttributeError, TypeError):
        pass

    # 언어 ID — python-pptx에 매핑 없는 로케일(en-IN 등) 시 ValueError 발생
    try:
        if font.language_id is not None:
            style["language_id"] = font.language_id
    except Exception:
        pass

    return style


def _style_key(style: dict) -> str:
    """스타일 딕셔너리를 비교 가능한 해시 키로 변환합니다."""
    items = sorted(style.items(), key=lambda x: x[0])
    return str(items)


def _extract_paragraph_style(paragraph) -> dict:
    """Paragraph의 서식 속성을 딕셔너리로 추출합니다."""
    pstyle: dict[str, Any] = {}
    if paragraph.alignment is not None:
        pstyle["alignment"] = paragraph.alignment
    if paragraph.level is not None:
        pstyle["level"] = paragraph.level
    try:
        if paragraph.line_spacing is not None:
            pstyle["line_spacing"] = paragraph.line_spacing
        if paragraph.space_before is not None:
            pstyle["space_before"] = paragraph.space_before
        if paragraph.space_after is not None:
            pstyle["space_after"] = paragraph.space_after
    except (AttributeError, TypeError):
        pass
    return pstyle


# ──────────────────────────────────────────────
#  슬라이드 맥락 추출
# ──────────────────────────────────────────────

def extract_slide_context(slide) -> str:
    """슬라이드의 모든 텍스트를 추출하여 맥락 문자열로 반환합니다."""
    texts = []
    for shape in slide.shapes:
        texts.extend(_extract_shape_texts(shape))
    return "\n".join(texts)


def _extract_shape_texts(shape) -> list[str]:
    """Shape에서 텍스트를 재귀적으로 추출합니다."""
    texts = []
    # 그룹 Shape 재귀 처리
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child_shape in shape.shapes:
            texts.extend(_extract_shape_texts(child_shape))
        return texts

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para_text = para.text.strip()
            if para_text:
                texts.append(para_text)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        texts.append(para_text)
    return texts


# ──────────────────────────────────────────────
#  스타일 ID 매핑 + 텍스트 추출
# ──────────────────────────────────────────────

def extract_styled_paragraphs(text_frame) -> dict | None:
    """
    TextFrame에서 paragraph/run 구조와 스타일 매핑을 추출합니다.

    Returns:
        {
            "styles": {"S0": {...}, "S1": {...}},
            "paragraphs": [
                {
                    "p_idx": 0,
                    "runs": [{"text": "Hello", "style_id": "S0"}, ...]
                },
                ...
            ]
        }
        텍스트가 없으면 None 반환.
    """
    styles_map: dict[str, dict] = {}
    key_to_id: dict[str, str] = {}
    style_counter = 0
    paragraphs_data = []
    has_text = False

    for p_idx, paragraph in enumerate(text_frame.paragraphs):
        runs_data = []
        for run in paragraph.runs:
            text = run.text
            if text:
                has_text = True

            rs = _extract_run_style(run)
            sk = _style_key(rs)

            if sk not in key_to_id:
                sid = f"S{style_counter}"
                style_counter += 1
                key_to_id[sk] = sid
                styles_map[sid] = rs

            runs_data.append({
                "text": text,
                "style_id": key_to_id[sk],
            })

        # Run이 없는 경우(순수 텍스트가 paragraph에 직접 있는 경우)
        if not runs_data and paragraph.text.strip():
            has_text = True
            runs_data.append({
                "text": paragraph.text,
                "style_id": "S0",
            })
            if "S0" not in styles_map:
                styles_map["S0"] = {}

        paragraphs_data.append({
            "p_idx": p_idx,
            "runs": runs_data,
        })

    if not has_text:
        return None

    return {
        "styles": styles_map,
        "paragraphs": paragraphs_data,
    }


# ──────────────────────────────────────────────
#  언어별 기본 폰트 매핑
# ──────────────────────────────────────────────

LANG_FONT_MAP = {
    # East Asian → <a:ea>
    "ko": "맑은 고딕",
    "ja": "游ゴシック",
    "zh": "Microsoft YaHei",
    # Complex Script → <a:cs>
    "ar": "Traditional Arabic",
    "th": "Leelawadee UI",
    # Cyrillic → <a:latin>
    "ru": "Calibri",
    # Vietnamese (확장 라틴) → <a:latin>
    "vi": "Segoe UI",
}

# 스크립트 유형별 언어 분류 (→ OOXML 폰트 요소 결정)
EAST_ASIAN_LANGS = {"ko", "ja", "zh"}
COMPLEX_SCRIPT_LANGS = {"ar", "th"}
LATIN_OVERRIDE_LANGS = {"ru", "vi"}


def _font_element_tag(target_lang: str) -> str | None:
    """대상 언어에 맞는 OOXML 폰트 요소 태그를 반환합니다."""
    if target_lang in EAST_ASIAN_LANGS:
        return "a:ea"
    if target_lang in COMPLEX_SCRIPT_LANGS:
        return "a:cs"
    if target_lang in LATIN_OVERRIDE_LANGS:
        return "a:latin"
    return None


def get_target_font(target_lang: str) -> str | None:
    """대상 언어에 맞는 기본 폰트를 반환합니다. 매핑이 없으면 None."""
    return LANG_FONT_MAP.get(target_lang)


# ──────────────────────────────────────────────
#  번역 결과 적용 (XML 직접 수정으로 서식 100% 보존)
# ──────────────────────────────────────────────


def _replace_run_text_xml(run, new_text):
    """
    Run의 텍스트를 XML 레벨에서 직접 교체합니다.
    <a:rPr> 등 서식 요소를 일체 건드리지 않으므로 서식이 100% 보존됩니다.
    """
    r_elem = run._r
    t_elem = r_elem.find(qn('a:t'))
    if t_elem is not None:
        t_elem.text = new_text
    else:
        # <a:t>가 없는 경우 새로 생성
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
        t_elem.text = new_text


def _set_run_target_font(run, target_font, target_lang):
    """
    Run에 대상 언어의 스크립트 유형에 맞는 폰트 요소를 설정합니다.
      - East Asian (ko/ja/zh)   → <a:ea>
      - Complex Script (ar/th)  → <a:cs>
      - Latin override (ru/vi)  → <a:latin>
    OOXML 스키마 순서(<a:latin> → <a:ea> → <a:cs>)를 준수합니다.
    기존 <a:rPr>의 다른 요소(색상, 크기 등)는 일체 건드리지 않습니다.
    """
    if not target_lang or not target_font:
        return

    tag = _font_element_tag(target_lang)
    if tag is None:
        return

    r_elem = run._r
    rPr = r_elem.find(qn('a:rPr'))
    if rPr is None:
        return

    # 이미 해당 요소가 있으면 typeface만 변경
    existing = rPr.find(qn(tag))
    if existing is not None:
        existing.set('typeface', target_font)
        return

    # 새 요소 생성
    new_elem = etree.Element(qn(tag))
    new_elem.set('typeface', target_font)

    # <a:latin>에서 pitchFamily, charset 속성 복사 (latin 자체가 아닌 경우)
    if tag != 'a:latin':
        latin = rPr.find(qn('a:latin'))
        if latin is not None:
            for attr in ('pitchFamily', 'charset'):
                val = latin.get(attr)
                if val is not None:
                    new_elem.set(attr, val)

    # OOXML 스키마 순서에 맞게 삽입: latin → ea → cs
    FONT_ORDER = ['a:latin', 'a:ea', 'a:cs']
    tag_idx = FONT_ORDER.index(tag)

    # 자신보다 뒤에 있는 요소를 찾아 그 앞에 삽입
    for later_tag in FONT_ORDER[tag_idx + 1:]:
        later = rPr.find(qn(later_tag))
        if later is not None:
            later.addprevious(new_elem)
            return

    # 자신보다 앞에 있는 요소를 찾아 그 뒤에 삽입
    for earlier_tag in reversed(FONT_ORDER[:tag_idx]):
        earlier = rPr.find(qn(earlier_tag))
        if earlier is not None:
            earlier.addnext(new_elem)
            return

    # 폰트 요소가 하나도 없으면 rPr 끝에 추가
    rPr.append(new_elem)


def apply_translated_runs(text_frame, translated_data: dict, styles_map: dict,
                          original_paragraphs_xml=None, styles_rPr_xml=None,
                          target_font: str | None = None,
                          target_lang: str | None = None):
    """
    번역된 데이터를 TextFrame에 적용합니다.

    ★ 핵심 원칙: XML 레벨에서 <a:t> 텍스트만 교체.
    기존 Run의 <a:rPr>(서식 정보)를 일체 건드리지 않아
    폰트, 색상, 크기, 볼드/이탤릭 등 서식을 100% 보존합니다.
    폰트는 PowerPoint 테마의 East Asian 폴백에 맡깁니다.

    전략:
    - paragraph 수가 같으면 1:1 매핑
    - 각 paragraph 내에서:
      - Run 수 동일 → 1:1 텍스트 교체 (최상의 스타일 보존)
      - Run 수 다름 → 첫 번째 Run에 전체 번역 텍스트, 나머지 비움

    Args:
        text_frame: python-pptx TextFrame 객체
        translated_data: {"paragraphs": [{"runs": [{"text": "...", "style_id": "S0"}, ...]}]}
        styles_map: 스타일 ID → 속성 매핑 (인터페이스 호환용)
        original_paragraphs_xml: (인터페이스 호환용)
        styles_rPr_xml: (인터페이스 호환용)
        target_font: (인터페이스 호환용, 현재 사용하지 않음)
        target_lang: (인터페이스 호환용, 현재 사용하지 않음)
    """
    translated_paras = translated_data.get("paragraphs", [])
    orig_paragraphs = list(text_frame.paragraphs)

    for p_idx, paragraph in enumerate(orig_paragraphs):
        orig_runs = list(paragraph.runs)
        if not orig_runs:
            continue  # Run이 없는 paragraph는 건드리지 않음

        # 대응하는 번역 paragraph가 있는지 확인
        if p_idx < len(translated_paras):
            t_para = translated_paras[p_idx]
            t_runs = t_para.get("runs", [])
            # 번역된 전체 텍스트
            full_translated = "".join(r.get("text", "") for r in t_runs)
        else:
            # 번역 paragraph가 부족하면 원문 유지
            continue

        if not full_translated and not any(r.get("text", "") for r in t_runs):
            # 번역 결과가 완전히 비어 있으면 원문 유지
            continue

        if len(t_runs) == len(orig_runs):
            # ── Run 수가 동일 → 1:1 매핑 (최상의 스타일 보존) ──
            for orig_run, t_run in zip(orig_runs, t_runs):
                new_text = t_run.get("text", "")
                _replace_run_text_xml(orig_run, new_text)
                if target_font and new_text:
                    _set_run_target_font(orig_run, target_font, target_lang)
            logger.debug(f"  P{p_idx}: 1:1 Run 매핑 ({len(orig_runs)}개)")
        else:
            # ── Run 수가 다름 → 첫 번째 Run에 전체 텍스트, 나머지 비움 ──
            _replace_run_text_xml(orig_runs[0], full_translated)
            if target_font and full_translated:
                _set_run_target_font(orig_runs[0], target_font, target_lang)
            for run in orig_runs[1:]:
                _replace_run_text_xml(run, "")
            logger.debug(
                f"  P{p_idx}: 통합 교체 (원본 {len(orig_runs)}개 Run → "
                f"번역 {len(t_runs)}개, 첫 Run에 통합)"
            )


# ──────────────────────────────────────────────
#  슬라이드 내 Shape 순회 (그룹 포함)
# ──────────────────────────────────────────────

def iter_translatable_shapes(slide):
    """
    슬라이드의 모든 번역 가능한 Shape를 반환합니다 (그룹 재귀 포함).
    각 항목은 (shape, "text_frame" | "table") 튜플.
    """
    yield from _iter_shapes_recursive(slide.shapes)


def _iter_shapes_recursive(shapes):
    """Shape 컬렉션을 재귀적으로 순회합니다."""
    for shape in shapes:
        # 그룹 Shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_shapes_recursive(shape.shapes)
            except AttributeError:
                pass
            continue

        if shape.has_text_frame:
            yield (shape, "text_frame")
        if shape.has_table:
            yield (shape, "table")


# ──────────────────────────────────────────────
#  원본 paragraph XML 백업/복원 유틸
# ──────────────────────────────────────────────

def backup_paragraph_xml(text_frame) -> list:
    """하위 호환용 — 새 방식에서는 사용하지 않습니다."""
    return []
